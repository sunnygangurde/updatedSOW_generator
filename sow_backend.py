# UPDATED CODE

import json
import uuid
import logging
import boto3
import os
import pdfplumber
from pptx import Presentation
from pptx.util import Pt
from docx import Document
from docx.shared import Pt as DocxPt
import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO
import re

# === Logging ===
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger()

# === AWS Clients ===
s3 = boto3.client("s3", region_name="us-east-1")
bedrock = boto3.client("bedrock-runtime", region_name="us-east-1")
BUCKET_NAME = 'sowbucketcreation'

# === Clean Input Text ===
def clean_proposal_text(raw_text):
    cleaned = re.sub(r'\n?\d{1,2}\s?[A-Z ]{3,}', '', raw_text)
    cleaned = re.sub(r'(Executive Summary|Our Understanding|Overall Scope of Work|Assumptions and Dependencies)', '', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r'\n+', '\n', cleaned)
    cleaned = re.sub(r'\s{2,}', ' ', cleaned)
    return cleaned.strip()

# === Remove Repeated Lines ===
def deduplicate_lines(text):
    seen = set()
    result = []
    for line in text.split("\n"):
        line = line.strip()
        if line and line not in seen:
            seen.add(line)
            result.append(line)
    return "\n".join(result)

# === Bedrock LLM Call ===
def call_bedrock(prompt):
    try:
        logger.info("Calling Bedrock with prompt...")
        response = bedrock.invoke_model(
            modelId="amazon.titan-text-express-v1",
            contentType="application/json",
            accept="application/json",
            body=json.dumps({
                "inputText": prompt,
                "textGenerationConfig": {
                    "maxTokenCount": 500,
                    "temperature": 0.3,
                    "topP": 0.9
                }
            })
        )
        result = json.loads(response["body"].read())
        return result["results"][0]["outputText"].strip()
    except Exception as e:
        logger.error(f"Bedrock error: {e}")
        return f"[ERROR generating section: {str(e)}]"

# === Section Generator ===
def generate_section(title, instruction, proposal_text, keywords=None, use_gdpr_template=False):
    if use_gdpr_template:
        try:
            with open("gdpr_appendix.txt", "r") as f:
                return f.read().strip()
        except Exception:
            return "Standard GDPR compliance policies will be added during the final contracting stage."

    if keywords and not any(k in proposal_text.lower() for k in keywords):
        return "To be defined during project discovery."

    prompt = f"""
You are a consultant. Write the **{title}** for an SoW. Instruction: {instruction}
Context:
{proposal_text[:2000]}
"""
    result = call_bedrock(prompt)
    return deduplicate_lines(result)

# === SoW Sections ===
SOW_SECTIONS = [
    ("DURATION", "Mention duration and expected timeline.", ["duration", "start", "end", "months", "weeks"]),
    ("SERVICES AND DELIVERABLES", "List services and deliverables.", ["deliverables", "scope", "services"]),
    ("IMPLEMENTATION TIMELINE", "Break into phases with time estimate.", ["timeline", "phase", "milestone"]),
    ("ACCEPTANCE CRITERIA", "List what defines project success.", ["acceptance", "criteria"]),
    ("GOVERNANCE AND MONITORING", "Mention reviews, stakeholders, issues.", ["governance", "monitoring"]),
    ("TEAM & ROLES", "List key roles and locations.", ["roles", "team"]),
    ("COMMERCIALS AND PAYMENT SCHEDULE", "Describe effort, milestones, and payments.", ["cost", "price"]),
    ("ASSUMPTIONS AND EXCLUSIONS", "Mention assumptions, dependencies, out-of-scope.", ["assumptions"]),
    ("DATA PROTECTION AND COMPLIANCE (e.g., GDPR)", "Mention data policies.", None),
    ("SIGN-OFF SECTION", "Include placeholders for sign-off.", ["signature"])
]

# === File Extractors ===
def extract_pdf(file_bytes):
    with open("temp_local.pdf", "wb") as f:
        f.write(file_bytes)
    with pdfplumber.open("temp_local.pdf") as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)

def extract_pptx(file_bytes):
    with open("temp_local.pptx", "wb") as f:
        f.write(file_bytes)
    prs = Presentation("temp_local.pptx")
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def extract_docx(file_bytes):
    with open("temp_local.docx", "wb") as f:
        f.write(file_bytes)
    doc = Document("temp_local.docx")
    return "\n".join(p.text for p in doc.paragraphs)

def extract_xlsx(file_bytes):
    with open("temp_local.xlsx", "wb") as f:
        f.write(file_bytes)
    df = pd.read_excel("temp_local.xlsx", sheet_name=None)
    return "\n\n".join([df[s].to_string(index=False) for s in df])

def extract_text(file_path):
    _, ext = os.path.splitext(file_path)
    with open(file_path, "rb") as f:
        file_bytes = f.read()

    ext = ext.lower()
    if ext == ".pdf":
        return extract_pdf(file_bytes)
    elif ext == ".pptx":
        return extract_pptx(file_bytes)
    elif ext == ".docx":
        return extract_docx(file_bytes)
    elif ext == ".xlsx":
        return extract_xlsx(file_bytes)
    else:
        return file_bytes.decode("utf-8", errors="ignore")

# === S3 Helpers ===
def upload_to_s3(key, data_bytes):
    s3.put_object(Bucket=BUCKET_NAME, Key=key, Body=data_bytes)
    logger.info(f"Uploaded to s3://{BUCKET_NAME}/{key}")

def generate_presigned_url(key, expiration=3600):
    return s3.generate_presigned_url('get_object', Params={'Bucket': BUCKET_NAME, 'Key': key}, ExpiresIn=expiration)

# === Exporters ===
def export_txt(content: str) -> bytes:
    return content.encode("utf-8")

def export_docx(content: str) -> bytes:
    doc = Document()
    doc.add_heading("Statement of Work for [Project Title]", 0)
    for section in content.split("### ")[1:]:
        if not section.strip():
            continue
        title_line, *body_lines = section.split("\n", 1)
        doc.add_heading(title_line.strip(), level=1)
        if body_lines:
            paragraphs = body_lines[0].strip().split("\n")
            for para in paragraphs:
                if para.strip():
                    p = doc.add_paragraph(para.strip())
                    run = p.runs[0]
                    run.font.size = DocxPt(11)
    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()

def export_pdf(content: str) -> bytes:
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    x, y = 50, height - 50
    for line in content.split("\n"):
        c.drawString(x, y, line[:90])
        y -= 15
        if y < 50:
            c.showPage()
            y = height - 50
    c.save()
    buffer.seek(0)
    return buffer.getvalue()

def export_pptx(content: str) -> bytes:
    prs = Presentation()
    for section in content.split("### ")[1:]:
        title, *body = section.split("\n", 1)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title.strip()
        slide.placeholders[1].text = body[0].strip() if body else ""
    output = BytesIO()
    prs.save(output)
    return output.getvalue()

def export_xlsx(content: str) -> bytes:
    df = pd.DataFrame([line.split(": ", 1) if ": " in line else [line, ""] for line in content.split("\n")])
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, header=["Section", "Content"])
    return output.getvalue()

# === Main Process ===
def process_file_and_generate_sow(input_file_path: str, output_format: str = "txt"):
    proposal_id = str(uuid.uuid4())
    raw_text = extract_text(input_file_path)
    proposal_text = clean_proposal_text(raw_text)

    if not proposal_text.strip():
        raise ValueError("No valid proposal text extracted.")

    section_outputs = []
    for title, instruction, keywords in SOW_SECTIONS:
        use_gdpr_template = "GDPR" in title.upper()
        content = generate_section(title, instruction, proposal_text, keywords, use_gdpr_template)
        section_outputs.append(f"### {title}\n\n{content.strip()}")

    final_sow = f"**Statement of Work for [Project Title]**\n\n" + "\n\n".join(section_outputs)

    if output_format == "pdf":
        sow_bytes = export_pdf(final_sow)
        file_key = f"sows/{proposal_id}.pdf"
    elif output_format == "docx":
        sow_bytes = export_docx(final_sow)
        file_key = f"sows/{proposal_id}.docx"
    elif output_format == "pptx":
        sow_bytes = export_pptx(final_sow)
        file_key = f"sows/{proposal_id}.pptx"
    elif output_format == "xlsx":
        sow_bytes = export_xlsx(final_sow)
        file_key = f"sows/{proposal_id}.xlsx"
    else:
        sow_bytes = export_txt(final_sow)
        file_key = f"sows/{proposal_id}.txt"

    upload_to_s3(file_key, sow_bytes)
    upload_to_s3(f'proposals/{proposal_id}.txt', proposal_text.encode("utf-8"))

    url = generate_presigned_url(file_key)
    return final_sow, url
